from pptx import Presentation
from pptx.util import Inches

from neuro_san.coded_tool import CodedTool


class PresentationGenerator(CodedTool):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)

    def create_presentation(self, filepath: str):
        """
        Creates a new PowerPoint presentation.

        :param filepath: The path to the new presentation.
        """
        prs = Presentation()
        prs.save(filepath)

    def add_slide(self, filepath: str, title: str, content: str):
        """
        Adds a new slide to an existing PowerPoint presentation.

        :param filepath: The path to the presentation.
        :param title: The title of the slide.
        :param content: The content of the slide.
        """
        prs = Presentation(filepath)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        title_shape.text = title
        tf = body_shape.text_frame
        tf.text = content
        prs.save(filepath)

    def add_picture(self, filepath: str, image_path: str, left: int, top: int, width: int, height: int):
        """
        Adds a picture to a slide in a PowerPoint presentation.

        :param filepath: The path to the presentation.
        :param image_path: The path to the image file.
        :param left: The distance from the left edge of the slide to the left edge of the picture.
        :param top: The distance from the top edge of the slide to the top edge of the picture.
        :param width: The width of the picture.
        :param height: The height of the picture.
        """
        prs = Presentation(filepath)
        slide = prs.slides[0]
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        prs.save(filepath)
